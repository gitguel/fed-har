#!/usr/bin/env python3
"""Monta os slides da orientação de 2026-08-05 sobre o template do Miguel.

Preserva o padrão do arquivo original (export do Google Slides): usa os MESMOS
layouts (`SECTION_HEADER` para divisor, `TITLE_AND_BODY_3` para conteúdo), a
mesma tipografia do corpo (negrito, hierarquia por `level`) e a mesma geometria
do slide 3 quando há figura à direita (corpo estreito, 7,4 de 20 pol).

Os 3 slides originais são preservados; tudo é acrescentado depois deles.
Conteúdo vem de `docs/apresentacao_05_08/roteiro_apresentacao_05_08.md`; figuras, de
`build_assets_05_08.py`.

Uso:
    poetry run python scripts/analysis/build_deck_05_08.py \
        --deck "docs/apresentacao_05_08/Apresentação 05_08.pptx" \
        --assets docs/apresentacao_05_08
"""
import argparse
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Emu, Inches, Pt

L_SECTION, L_BODY = 2, 5          # índices dos layouts no master
# Geometria do corpo estreito do slide 3 — deixa a direita livre para figura.
BODY_NARROW = dict(left=Inches(0.42), top=Inches(3.50),
                   width=Inches(7.40), height=Inches(6.99))
FIG_BOX = dict(left=Inches(8.20), top=Inches(2.60),
               width=Inches(11.30), height=Inches(7.90))


def _set_body(ph, itens, size):
    """Escreve bullets com hierarquia. `itens` = [(nivel, texto), ...]."""
    tf = ph.text_frame
    tf.clear()
    for i, (lvl, txt) in enumerate(itens):
        par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        par.level = lvl
        run = par.add_run()
        run.text = txt
        run.font.size = Pt(size)
        run.font.bold = True          # o corpo do template é negrito


def _fit(box, w, h):
    """Encaixa (w,h) na caixa preservando proporção; devolve (l, t, w, h) em EMU."""
    esc = min(box["width"] / w, box["height"] / h)
    nw, nh = int(w * esc), int(h * esc)
    return (int(box["left"] + (box["width"] - nw) / 2),
            int(box["top"] + (box["height"] - nh) / 2), nw, nh)


class Deck:
    def __init__(self, path):
        self.p = Presentation(path)
        self.layouts = self.p.slide_masters[0].slide_layouts

    def secao(self, titulo):
        s = self.p.slides.add_slide(self.layouts[L_SECTION])
        s.placeholders[0].text_frame.text = titulo
        for ph in list(s.placeholders):          # subtítulo vazio: remover
            if ph.placeholder_format.idx == 1 and not ph.text_frame.text:
                ph._element.getparent().remove(ph._element)
        return s

    def conteudo(self, titulo, itens, sub="", figura=None, size=24):
        s = self.p.slides.add_slide(self.layouts[L_BODY])
        alvo = {ph.placeholder_format.idx: ph for ph in s.placeholders}
        alvo[0].text_frame.text = titulo
        if 3 in alvo:
            if sub:
                tf = alvo[3].text_frame
                tf.text = sub
                for par in tf.paragraphs:
                    for r in par.runs:
                        r.font.size = Pt(16)
            else:
                alvo[3]._element.getparent().remove(alvo[3]._element)
        if 2 in alvo:                             # faixa extra do layout, não usada
            alvo[2]._element.getparent().remove(alvo[2]._element)
        body = alvo[1]
        _set_body(body, itens, size)
        if figura:
            for k, v in BODY_NARROW.items():
                setattr(body, k, v)
            with Image.open(figura) as im:
                w, h = im.size
            l, t, fw, fh = _fit(FIG_BOX, w, h)
            s.shapes.add_picture(str(figura), l, t, fw, fh)
        return s


def main():
    ap = argparse.ArgumentParser(description=__doc__.split("\n")[0])
    ap.add_argument("--deck", required=True)
    ap.add_argument("--assets", required=True)
    a = ap.parse_args()
    A = Path(a.assets)
    d = Deck(a.deck)

    # ---------------------------------------------- Bloco 1: o que foi rodado
    d.conteudo("Setup Federado", [
        (0, "spec = quem treina"), (1, "a composição de clientes da federação"),
        (0, "target = onde testa"), (1, "o test.csv que avalia o modelo"),
        (0, "shots = quanto rótulo"), (1, "por classe, por cliente, no fine-tuning"),
    ], sub="Três palavras que precisam estar claras antes de qualquer número", size=26)

    d.conteudo("4 federações → 6 avaliações", [
        (0, "in10-RW: 10 usuários RealWorld_thigh → testa em RW"),
        (0, "in10-MS: 10 usuários MotionSense → testa em MS"),
        (0, "cross5+5: 5 RW + 5 MS → testa em RW e MS"),
        (0, "cross10+10: 10 RW + 10 MS → testa em RW e MS"),
        (0, ""),
        (0, "Um cross é UMA federação avaliada em dois testes"),
        (1, "mesmo modelo, dois números — não são dois treinos"),
    ], sub="4 treinos, 6 células de avaliação", size=24)

    d.conteudo("Duas fases de treino", [
        (0, "Fase 1 — pré-treino SSL (sem rótulo)"),
        (1, "R = 100 rodadas · TF-C 5 épocas locais, LFR 30"),
        (1, "produz o backbone"),
        (0, "Fase 2 — fine-tuning supervisionado"),
        (1, "R = 150 rodadas · 5 épocas locais"),
        (1, "parte do backbone da Fase 1"),
        (0, ""),
        (0, "Um backbone serve os 5 degraus de rótulo"),
        (1, "o pré-treino não usa rótulo — o orçamento não o afeta"),
    ], sub="O baseline pula a Fase 1: começa aleatório", size=22)

    d.conteudo("O tamanho da grade", [
        (0, "192 pré-treinos federados"),
        (0, "1.632 fine-tunings federados"),
        (0, "2.112 resultados (run × alvo)"),
        (1, "um run rende 1 ou 2, conforme o spec tenha 1 ou 2 domínios"),
        (0, ""),
        (0, "Eixos: 4 encoders × 3 métodos × 4 specs"),
        (1, "× 5 degraus de rótulo × 4 seeds × 2 domínios"),
    ], sub="Em execução: mais 64 pré-treinos e 320 fine-tunings (Exp. 2 @full)",
        size=24)

    d.conteudo("Avaliamos o modelo global", [
        (0, "Toda avaliação é do modelo agregado por FedAvg"),
        (1, "não há modelo local nem personalização nesta grade"),
        (0, ""),
        (0, "Teste em sujeitos NUNCA vistos"),
        (1, "RW: treina em 10 usuários, testa em 3 outros"),
        (1, "MS: treina em 17, testa em 5 outros"),
        (1, "interseção train ∩ test = zero"),
    ], sub="Splits do DAGHAR são por usuário — sem vazamento por sujeito", size=24)

    # ---------------------------------------------- Bloco 2: hiperparâmetros
    d.secao("Hiperparâmetros")
    d.conteudo("Por que cada escolha", [
        (0, "batch 64 — benchmark do da Luz (Seção V-D)"),
        (0, "192 janelas/cliente = 3 batches"),
        (1, "mínimo que dá ao TF-C mais de um passo por época"),
        (0, "lr 1e-4 no fine-tuning — melhor tunado (Tabela 12)"),
        (0, "Adam + CrossEntropy — default do benchmark"),
        (0, "5 épocas locais — consenso de 4 papers primários"),
        (0, "LFR 30 = 6×5 — casa épocas efetivas de backbone"),
    ], sub="Tudo herdado do benchmark, para comparabilidade célula-a-célula", size=22)

    d.conteudo("Uma ressalva que prefiro dizer antes", [
        (0, "lr do fine-tuning (1e-4): melhor tunado do benchmark"),
        (0, "lr do pré-treino (3e-4): default do YAML"),
        (1, "idêntico nos 2 métodos e nos 4 encoders"),
        (1, "não existe varredura publicada — nem nossa, nem do benchmark"),
        (0, ""),
        (0, "→ vira item de trabalho futuro"),
    ], sub="Os hiperparâmetros de pré-treino não são “os melhores”", size=24)

    # ---------------------------------------------- Bloco 3: resultados
    d.secao("Resultados")

    d.conteudo("H1 — domain shift prejudica a federação", [
        (0, "Dois contrastes, duas forças:"),
        (1, "diluição (5+5 − in10): −5,5 pp, 7 de 8 células"),
        (1, "aumento (10+10 − in10): −1,8 pp, 6 de 8"),
        (0, ""),
        (0, "A hipótese-pilar se sustenta"),
        (1, "e a queda acompanha a distância entre domínios"),
        (1, "medida em runs centralizados independentes"),
    ], sub="Se ela não valesse, não haveria problema para o Fed-SSL atacar",
        figura=A / "fig_h1_contrastes.png", size=21)

    d.conteudo("H1 — o custo cresce com o rótulo", [
        (0, "diluição: −3,0 → −5,5 pp (L=1 → full)"),
        (0, "aumento: +1,1 → −1,8 pp"),
        (0, ""),
        (0, "Com 1 rótulo por classe, dado estrangeiro AJUDA"),
        (0, ""),
        (0, "Contraria a motivação usual do SSL"),
        (1, "que supõe o shift machucar mais quando falta rótulo"),
        (1, "aqui o regime few-shot é onde ele menos machuca"),
    ], sub="O oposto da intuição — e muda como vendemos o Fed-SSL", size=23)

    d.conteudo("H2 — o efeito depende do par (método, encoder)", [
        (0, "Não existe “o efeito do SSL”"),
        (0, ""),
        (0, "tfc + rnn: +20,3 pp, 100% das 40 células"),
        (0, "tfc + tstcc: −2,1 pp (−4,9 no full)"),
        (0, "os outros seis: indistinguíveis entre si"),
        (0, ""),
        (0, "TF-C resgata o PIOR encoder da grade"),
        (1, "não é “o TF-C é ótimo”"),
    ], sub="rnn — o encoder onde o TF-C mais entrega",
        figura=A / "fig_H2_rnn.png", size=21)

    d.conteudo("H2 — e onde ele atrapalha", [
        (0, "tstcc: TF-C fica ABAIXO do supervisionado"),
        (1, "sempre que há domínio estrangeiro"),
        (0, ""),
        (0, "MS: 80,8 → 73,8 (cross5+5)"),
        (0, "RW: 65,8 → 63,0"),
        (0, ""),
        (0, "Enquanto o LFR sobe no mesmo encoder"),
    ], sub="tstcc — o sinal oposto, mesmo método",
        figura=A / "fig_H2_tstcc.png", size=21)

    d.conteudo("O teste: é estrutura ou é ruído?", [
        (0, "Variância entre pares ÷ ruído de seed:"),
        (0, ""),
        (0, "Δ A, grade completa: 2,21"),
        (0, "Δ A, sem tfc+rnn: 0,16"),
        (0, "DiD: 0,11"),
        (0, ""),
        (0, "A heterogeneidade é INTEIRAMENTE o tfc + rnn"),
        (1, "ordenar os outros sete é ordenar ruído"),
    ], sub="Com 4 seeds, o dp é 4,7 pp — diferenças menores não se distinguem", size=23)

    d.conteudo("H3 — o SSL não ataca o shift", [
        (0, "DiD ≈ 0 ou negativo em toda a grade"),
        (1, "razão 0,11 com ou sem o outlier"),
        (0, ""),
        (0, "Δ A é a pergunta de deployment"),
        (1, "“vou operar cross-domain — o SSL ajuda?” → sim, em alguns pares"),
        (0, "DiD é a pergunta mecanística"),
        (1, "“o SSL ataca o shift?” → não"),
        (0, ""),
        (0, "Resultado negativo, não hipótese fracassada"),
    ], sub="A motivação original do projeto afirma que sim — e os dados dizem que não",
        size=21)

    d.conteudo("Exp. 2 — quanto custa federar o pré-treino", [
        (0, "Com dado fixo, federar custa ~1 pp"),
        (1, "pequeno perto do +7,9 pp que o SSL entrega no 1-shot"),
        (0, ""),
        (0, "E é assimétrico:"),
        (1, "TF-C paga ~0"),
        (1, "LFR paga ~2 pp"),
        (0, ""),
        (0, "→ a perda é da loss batch-hungry, não do protocolo"),
        (0, "Braço @full rodando agora (separa volume de federação)"),
    ], sub="320 células por braço, fechado hoje de manhã", size=22)

    # ---------------------------------------------- Bloco 4 e 5
    d.secao("Próximos passos")

    d.conteudo("Fechar este setup", [
        (0, "1. FedSimCLR — prioridade máxima"),
        (1, "baseline mais usado na literatura de FedSSL"),
        (1, "já está na minerva; fecha o triângulo das 3 famílias"),
        (1, "é batch-hungry → toca o achado do piso de batch"),
        (0, "2. TS2Vec e encoders restantes do benchmark"),
        (0, "3. Avaliação out-of-domain no eixo federado"),
        (1, "custo quase zero: 1.632 checkpoints salvos, --targets já existe"),
        (0, "4. Varredura de lr do pré-treino"),
    ], sub="Em ordem do que um revisor cobra primeiro", size=21)

    d.conteudo("Trabalhos futuros", [
        (0, "Piso de batch em FSSL — contribuição nomeada"),
        (1, "KuHar: 48 de 57 clientes (84%) têm menos que um batch"),
        (1, "o cliente mínimo viável em FSSL é um batch, não uma amostra"),
        (0, "Personalização — começa sem re-rodar nada"),
        (1, "RW: 832 janelas/cliente ociosas, nunca vistas"),
        (0, "Mais datasets: 2ª onda cintura↔perna (gap maior)"),
        (0, "Ablações: F7, backbone-only, FedBN-SSL, agg_shock"),
    ], sub="Coletados durante a implementação — os dados ou o código já apontaram",
        size=21)

    d.p.save(a.deck)
    print(f"[DECK] {len(d.p.slides.__iter__.__self__._sldIdLst)} slides -> {a.deck}")


if __name__ == "__main__":
    main()
